import io, datetime
import pandas as pd
import plotly.io as pio
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from pptx import Presentation
from pptx.util import Inches

def export_excel(df_filtered: pd.DataFrame) -> bytes:
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="xlsxwriter") as xw:
        df_filtered.to_excel(xw, sheet_name="Detailed", index=False)

        summary = (
            df_filtered
            .groupby(["Accounting_Period","Type"], as_index=False)[["Budget","Actual"]]
            .sum()
        )
        summary["Variance_£"] = summary["Actual"] - summary["Budget"]
        summary["Variance_%"] = summary["Variance_£"] / summary["Budget"].replace(0, float("nan"))
        summary.to_excel(xw, sheet_name="Summary", index=False)
    out.seek(0)
    return out.read()

def export_pdf(figs: list) -> bytes:
    out = io.BytesIO()
    c = canvas.Canvas(out, pagesize=A4)
    width, height = A4

    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, height-50, "FinSight AI — Variance Summary")
    c.setFont("Helvetica", 10)
    c.drawString(40, height-65, datetime.datetime.now().strftime("%Y-%m-%d %H:%M"))

    y = height - 120
    for fig in figs:
        img_bytes = pio.to_image(fig, format="png", scale=2)
        c.drawInlineImage(io.BytesIO(img_bytes), 40, max(80, y-320), width=515, height=300)
        y -= 340
        if y < 120:
            c.showPage()
            y = height - 120

    c.save()
    out.seek(0)
    return out.read()

def export_pptx(figs: list) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FinSight AI — Variance Dashboard"
    slide.placeholders[1].text = "Auto-generated board slides"

    for fig in figs:
        img = pio.to_image(fig, format="png", scale=2)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        left = Inches(0.5); top = Inches(0.5); width = Inches(12.3)
        slide.shapes.add_picture(io.BytesIO(img), left, top, width=width)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
