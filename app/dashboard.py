# app/dashboard.py
import io
from pathlib import Path
import pandas as pd
import numpy as np
import streamlit as st
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
pio.templates.default = "plotly_white"

from utils.chatbot.engine import answer_query, append_history  # Chatbot integration

# --------------------------- Brand & Layout ---------------------------
BRAND = {
    "name": "FinSight AI",
    "colors": {
        "primary": "#B9D8B3",  # sage
        "dark": "#2C3E2F",  # forest
        "cream": "#FFF5E6",  # background
        "charcoal": "#3A3D3A",  # text
        "red": "#D84C4C",
        "green": "#9CCB8F",
    },
}

st.set_page_config(page_title="FinSight AI — Variance Demo", layout="wide", page_icon="📊")

st.markdown(
    f"""
<style>
    .stApp {{
        background-color: {BRAND["colors"]["cream"]};
    }}
    .metric .metric-value {{
        color: {BRAND["colors"]["dark"]} !important;
    }}
    section[data-testid="stSidebar"] {{
        width: 320px;
    }}
</style>
""",
    unsafe_allow_html=True,
)

# --------------------------- Data Load ---------------------------
DATA_PATH = Path(__file__).parent / "data.csv"


@st.cache_data
def load_data(csv_path: Path) -> pd.DataFrame:
    df = pd.read_csv(csv_path)

    for col in ["Accounting_Period", "Department", "Category", "Product"]:
        if col not in df.columns:
            df[col] = "All"

    for col in ["Revenue", "Direct_Costs", "Indirect_Costs", "Actual", "Budget"]:
        if col not in df.columns:
            df[col] = 0.0

    if "Variance_£" not in df.columns:
        df["Variance_£"] = df["Actual"] - df["Budget"]
    if "Variance_%" not in df.columns:
        df["Variance_%"] = np.where(
            df["Budget"].replace(0, np.nan).notna(), df["Variance_£"] / df["Budget"], 0.0
        )
    return df


df_raw = load_data(DATA_PATH)

# --------------------------- Sidebar Filters ---------------------------
st.sidebar.header("Data")
st.sidebar.caption("Upload CSV (optional)")
upload = st.sidebar.file_uploader(" ", type=["csv"], label_visibility="collapsed")
if upload is not None:
    tmp = pd.read_csv(upload)
    df_raw = load_data(Path(upload.name)) if False else tmp

st.sidebar.header("Filters")


def ms(label, series: pd.Series):
    opts = sorted(series.dropna().unique().tolist())
    return st.sidebar.multiselect(label, opts, default=opts)


sel_periods = ms("Accounting Period", df_raw["Accounting_Period"])
sel_departments = ms("Department", df_raw["Department"])
sel_categories = ms("Category", df_raw["Category"])
sel_products = ms("Product", df_raw["Product"])

mask = (
    df_raw["Accounting_Period"].isin(sel_periods)
    & df_raw["Department"].isin(sel_departments)
    & df_raw["Category"].isin(sel_categories)
    & df_raw["Product"].isin(sel_products)
)
df = df_raw.loc[mask].copy()

# --------------------------- Helpers ---------------------------
def money(x: float) -> str:
    try:
        return f"£{x:,.0f}"
    except Exception:
        return str(x)


def kpi_tiles_financials(df_: pd.DataFrame):
    revenue = float(df_["Revenue"].sum())
    direct = float(df_["Direct_Costs"].sum())
    indirect = float(df_["Indirect_Costs"].sum())
    ebitda = revenue - direct - indirect

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        st.metric("Revenue", money(revenue))
    with c2:
        st.metric("Direct Costs", money(direct))
    with c3:
        st.metric("Indirect Costs", money(indirect))
    with c4:
        st.metric("EBITDA", money(ebitda))


def executive_summary_bullets(df_: pd.DataFrame) -> list[str]:
    bullets = []
    total_actual = float(df_["Actual"].sum())
    total_budget = float(df_["Budget"].sum())
    var = total_actual - total_budget
    pct = (var / total_budget) if total_budget else 0.0

    direction = "over" if var > 0 else "under"
    bullets.append(
        f"Actual is **{direction}** budget by **{money(abs(var))} ({pct:.1%})** overall."
    )

    if "Revenue" in df_.columns and "Category" in df_.columns:
        rev = df_.groupby("Category", as_index=False)["Revenue"].sum()
        if len(rev):
            top = rev.sort_values("Revenue", ascending=False).iloc[0]
            bullets.append(f"Top revenue driver: **{top['Category']}** at {money(top['Revenue'])}.")

    freight = df_[df_["Category"].str.contains("Freight", case=False, na=False)]
    if len(freight):
        fv = float(freight["Actual"].sum() - freight["Budget"].sum())
        if abs(fv) > 0:
            label = "over" if fv > 0 else "under"
            bullets.append(f"Freight is **{label}** budget by {money(abs(fv))}.")

    if not bullets:
        bullets.append("No strong signals detected. Try widening filters.")
    return bullets


# --------------------------- Title ---------------------------
st.title("FinSight AI — Variance Analysis (Demo)")
st.caption("Minimalist, board-ready visuals • Sidebar filters are the single source of truth")

# --------------------------- KPI ROWS ---------------------------
kpi_tiles_financials(df)

tot_actual = float(df["Actual"].sum())
tot_budget = float(df["Budget"].sum())
tot_var = tot_actual - tot_budget
tot_varpct = (tot_var / tot_budget) if tot_budget else 0.0

k1, k2, k3, k4 = st.columns(4)
with k1:
    st.metric("Actual", money(tot_actual))
with k2:
    st.metric("Budget", money(tot_budget))
with k3:
    st.metric("Variance £", money(tot_var))
with k4:
    st.metric("Variance %", f"{tot_varpct:.1%}")

# --------------------------- Executive Summary ---------------------------
st.subheader("Executive Summary 🤖")
for line in executive_summary_bullets(df):
    st.markdown(f"- {line}")

# --------------------------- Variance by Period (bar) ---------------------------
st.subheader("Variance Overview")
if {"Actual", "Budget", "Accounting_Period"}.issubset(df.columns):
    g = df.copy()
    g["Variance_£"] = g["Actual"] - g["Budget"]
    per = (
        g.groupby("Accounting_Period", as_index=False)["Variance_£"]
        .sum()
        .sort_values("Accounting_Period")
    )
    fig_bar = px.bar(
        per,
        x="Accounting_Period",
        y="Variance_£",
        color="Variance_£",
        color_continuous_scale=[BRAND["colors"]["red"], BRAND["colors"]["green"]],
        labels={"Variance_£": "Variance (£)"},
    )
    fig_bar.update_layout(coloraxis_showscale=False, plot_bgcolor="rgba(0,0,0,0)")
    st.plotly_chart(fig_bar, use_container_width=True)
else:
    st.info("Need Actual, Budget, and Accounting_Period for the variance bar.")

# --------------------------- Variance % Trend ---------------------------
st.subheader("Variance % Trend")
if {"Actual", "Budget", "Accounting_Period"}.issubset(df.columns):
    g = df.copy()
    g["Variance_%"] = np.where(
        g["Budget"] != 0, (g["Actual"] - g["Budget"]) / g["Budget"], 0.0
    )
    per = (
        g.groupby("Accounting_Period", as_index=False)["Variance_%"]
        .mean()
        .sort_values("Accounting_Period")
    )
    fig_line = px.line(
        per,
        x="Accounting_Period",
        y="Variance_%",
        markers=True,
        labels={"Variance_%": "Variance %"},
    )
    fig_line.update_layout(plot_bgcolor="rgba(0,0,0,0)")
    st.plotly_chart(fig_line, use_container_width=True)
else:
    st.info("Need Actual, Budget, and Accounting_Period for the variance % line.")

# --------------------------- Waterfall ---------------------------
st.subheader("Waterfall — Budget to Actual (by Category impact)")
if {"Actual", "Budget", "Category"}.issubset(df.columns):
    total_budget = float(df["Budget"].sum())
    total_actual = float(df["Actual"].sum())
    cat_var = (
        df.assign(Var=lambda d: d["Actual"] - d["Budget"])
        .groupby("Category", as_index=False)["Var"]
        .sum()
        .sort_values("Var", ascending=True)
    )

    x_labels = ["Total Budget"] + cat_var["Category"].tolist() + ["Total Actual"]
    measures = ["absolute"] + ["relative"] * len(cat_var) + ["total"]
    y_values = [total_budget] + cat_var["Var"].tolist() + [total_actual]

    fig_wf = go.Figure(
        go.Waterfall(
            x=x_labels,
            measure=measures,
            y=y_values,
            connector={"line": {"color": BRAND["colors"]["charcoal"]}},
            increasing={"marker": {"color": BRAND["colors"]["green"]}},
            decreasing={"marker": {"color": BRAND["colors"]["red"]}},
            totals={"marker": {"color": BRAND["colors"]["dark"]}},
        )
    )
    fig_wf.update_layout(showlegend=False, plot_bgcolor="rgba(0,0,0,0)")
    st.plotly_chart(fig_wf, use_container_width=True)
else:
    st.info("Need Actual, Budget, and Category columns for the waterfall.")

# --------------------------- Detailed Variance Table ---------------------------
st.subheader("Detailed Variance Table")
detail_cols = [
    c
    for c in [
        "Accounting_Period",
        "Department",
        "Category",
        "Product",
        "Budget",
        "Actual",
        "Variance_£",
        "Variance_%",
    ]
    if c in df.columns
]
table_df = df[detail_cols] if detail_cols else df
st.dataframe(
    table_df.sort_values(
        detail_cols[:4] if len(detail_cols) >= 4 else table_df.columns.tolist()
    ),
    use_container_width=True,
)

# --------------------------- Export Options ---------------------------
st.subheader("Export Options")

out_xlsx = io.BytesIO()
with pd.ExcelWriter(out_xlsx, engine="xlsxwriter") as writer:
    table_df.to_excel(writer, index=False, sheet_name="Variance Detail")
st.download_button(
    "📥 Export to Excel",
    data=out_xlsx.getvalue(),
    file_name="finsight_variance_detail.xlsx",
    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
)

try:
    from pptx import Presentation
    from pptx.util import Inches

    def build_ppt(df_table: pd.DataFrame) -> bytes:
        prs = Presentation()
        title = prs.slides.add_slide(prs.slide_layouts[0])
        title.shapes.title.text = "FinSight AI — Variance Summary"
        title.placeholders[1].text = "Auto-generated export for board-ready review."

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        rows = min(18, len(df_table) + 1)
        cols = min(len(df_table.columns), 8)
        table = shapes.add_table(
            rows, cols, Inches(0.5), Inches(1.0), Inches(12.8), Inches(5.5)
        ).table

        headers = list(df_table.columns)[:cols]
        for j, h in enumerate(headers):
            table.cell(0, j).text = str(h)
        for i in range(1, rows):
            if i - 1 >= len(df_table):
                break
            for j in range(cols):
                table.cell(i, j).text = str(df_table.iloc[i - 1, j])

        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()

    ppt_bytes = build_ppt(table_df)
    st.download_button(
        "🖼️ Export to PowerPoint",
        data=ppt_bytes,
        file_name="finsight_variance.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    st.download_button(
        "🟡 Export to Google Slides",
        data=ppt_bytes,
        file_name="finsight_variance_for_slides.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        help="Import this .pptx into Google Slides.",
    )
except Exception:
    st.warning("PowerPoint export requires `python-pptx` (pip install python-pptx).")

# --------------------------- Chatbot ---------------------------
st.subheader("Ask your data 💬")
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

for msg in st.session_state.chat_history:
    role = "🧑‍💼" if msg["role"] == "user" else "🤖"
    st.markdown(f"**{role} {msg['role'].title()}:** {msg['content']}")

q = st.text_input("Type a question (e.g., 'Top 5 overspends', 'EBITDA in Jun 25')", "")
if st.button("Ask") and q.strip():
    ctx = {
        "Product": sel_products,
        "Department": sel_departments,
        "Category": sel_categories,
        "Accounting_Period": sel_periods,
    }
    st.session_state.chat_history = append_history(st.session_state.chat_history, "user", q)
    reply = answer_query(df, q, ctx)
    st.session_state.chat_history = append_history(st.session_state.chat_history, "assistant", reply)
    st.experimental_rerun()
